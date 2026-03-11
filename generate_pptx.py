"""Generate the 4-slide executive summary PowerPoint for Apple Patent Explorer."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
BG = RGBColor(0x1A, 0x1B, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DIMMED = RGBColor(0x8B, 0x92, 0x9E)
GREEN = RGBColor(0x51, 0xCF, 0x66)
BLUE = RGBColor(0x33, 0x9A, 0xF0)
RED = RGBColor(0xFF, 0x6B, 0x6B)
ORANGE = RGBColor(0xFF, 0x92, 0x2B)
YELLOW = RGBColor(0xFC, 0xC4, 0x19)
CYAN = RGBColor(0x22, 0xB8, 0xCF)
PURPLE = RGBColor(0xCC, 0x5D, 0xE8)
CARD_BG = RGBColor(0x25, 0x26, 0x2B)
BORDER = RGBColor(0x2C, 0x2E, 0x33)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_text(slide, left, top, width, height, text, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_card(slide, left, top, width, height, label, value, color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)
    # Label
    add_text(slide, left + 0.15, top + 0.1, width - 0.3, 0.3, label, size=9, color=DIMMED, bold=True)
    # Value
    add_text(slide, left + 0.15, top + 0.4, width - 0.3, 0.5, str(value), size=22, color=color, bold=True)

def add_bar(slide, left, top, width, max_width, height, color, label, value):
    """Draw a simple horizontal bar with label."""
    # Bar background
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(max_width), Inches(height))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = RGBColor(0x2A, 0x2E, 0x35)
    bg_shape.line.fill.background()
    # Bar fill
    if width > 0.05:
        fill_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        fill_shape.fill.solid()
        fill_shape.fill.fore_color.rgb = color
        fill_shape.line.fill.background()
    # Label
    add_text(slide, left + max_width + 0.1, top - 0.03, 3, height + 0.06, f"{label}: {value}", size=10, color=DIMMED)


# ============================================================
# SLIDE 1: Executive Summary
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide1)

add_text(slide1, 0.5, 0.3, 12, 0.6, "Apple Patent Explorer", size=32, color=WHITE, bold=True)
add_text(slide1, 0.5, 0.85, 12, 0.4, "Executive Summary  |  Design Patent Portfolio Analytics", size=16, color=DIMMED)

# Divider line
shape = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.35), Inches(12.3), Pt(2))
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE
shape.line.fill.background()

# What this is
add_text(slide1, 0.5, 1.6, 6, 0.3, "WHAT IS THIS?", size=11, color=BLUE, bold=True)
add_text(slide1, 0.5, 1.9, 6, 1.4,
    "A custom-built web application that transforms 3,000 Apple design patent records "
    "from the USPTO (U.S. Patent Office) into an interactive dashboard. Instead of "
    "scrolling through government databases or expensive legal software, this tool "
    "gives you instant charts, tables, and numbers to understand Apple's design "
    "patent portfolio at a glance.",
    size=12, color=WHITE)

add_text(slide1, 0.5, 3.2, 6, 0.3, "WHERE DOES THE DATA COME FROM?", size=11, color=BLUE, bold=True)
add_text(slide1, 0.5, 3.5, 6, 1.2,
    "Raw patent data is pulled from the USPTO Patent Examination Data System (PEDS) "
    "API — the same source patent attorneys and IP professionals use. The raw data "
    "was first cleaned using a Python script (clean_apple.py) that parses inventor "
    "names into structured fields, standardizes titles, and validates date formats. "
    "The cleaned data then feeds into the dashboard.",
    size=12, color=WHITE)

# KPI cards on right side
add_text(slide1, 7.3, 1.6, 5, 0.3, "KEY METRICS", size=11, color=BLUE, bold=True)
add_card(slide1, 7.3, 2.0, 2.4, 0.9, "TOTAL APPLICATIONS", "3,000", WHITE)
add_card(slide1, 10.0, 2.0, 2.4, 0.9, "GRANTED", "2,961", GREEN)
add_card(slide1, 7.3, 3.1, 2.4, 0.9, "GRANT RATE", "98.7%", CYAN)
add_card(slide1, 10.0, 3.1, 2.4, 0.9, "AVG TIME TO GRANT", "19.2 mo", YELLOW)
add_card(slide1, 7.3, 4.2, 2.4, 0.9, "UNIQUE INVENTORS", "1,680", PURPLE)
add_card(slide1, 10.0, 4.2, 2.4, 0.9, "ABANDONED", "29", RED)
add_card(slide1, 7.3, 5.3, 2.4, 0.9, "PENDING / ACTIVE", "10", BLUE)
add_card(slide1, 10.0, 5.3, 2.4, 0.9, "PORTFOLIO TYPE", "100% Design", ORANGE)

# Bottom summary
add_text(slide1, 0.5, 4.8, 6, 0.3, "BOTTOM LINE", size=11, color=GREEN, bold=True)
add_text(slide1, 0.5, 5.1, 6, 1.8,
    "Apple's design patent portfolio is remarkably healthy. A 98.7% grant rate with an "
    "average 19-month turnaround is well above industry averages. Only 29 out of 3,000 "
    "applications were abandoned. The portfolio spans 2018-2025 with peak filing in 2022 "
    "(502 applications). This dashboard replaces manual spreadsheet tracking and expensive "
    "third-party patent analytics platforms.",
    size=12, color=WHITE)

add_text(slide1, 0.5, 7.0, 12, 0.3, "Apple Inc.  |  Design Patent Portfolio  |  3,000 Applications  |  2018-2025", size=10, color=DIMMED, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: Prosecution & Status Analysis (OAs, statuses, abandonments)
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide2)

add_text(slide2, 0.5, 0.3, 12, 0.6, "Prosecution & Status Analysis", size=28, color=WHITE, bold=True)
add_text(slide2, 0.5, 0.8, 12, 0.3, "Office Actions, Application Statuses, Abandonments & Pipeline", size=14, color=DIMMED)

shape = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.3), Pt(2))
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE
shape.line.fill.background()

# LEFT: Status breakdown
add_text(slide2, 0.5, 1.45, 5, 0.3, "APPLICATION STATUS BREAKDOWN", size=11, color=BLUE, bold=True)

statuses = [
    ("Patented Case", 2961, GREEN),
    ("Abandoned (Failure to Respond to OA)", 25, RED),
    ("Notice of Allowance Mailed", 3, CYAN),
    ("Abandoned (Failure to Pay Issue Fee)", 3, RED),
    ("Docketed / Ready for Examination", 2, BLUE),
    ("Response to Non-Final OA Entered", 2, ORANGE),
    ("Abandoned (After Examiner's Answer)", 1, RED),
    ("Non-Final Action Mailed", 1, ORANGE),
]

y = 1.85
max_bar = 5.5
max_val = 2961
for label, val, color in statuses:
    bar_w = max(0.05, (val / max_val) * max_bar)
    add_bar(slide2, 0.5, y, bar_w, max_bar, 0.22, color, label, f"{val:,}")
    y += 0.42

# RIGHT: Prosecution insights
add_text(slide2, 7.3, 1.45, 5.5, 0.3, "PROSECUTION INSIGHTS", size=11, color=BLUE, bold=True)

# Cards for key prosecution metrics
add_card(slide2, 7.3, 1.85, 2.5, 0.85, "TOTAL OFFICE ACTIONS", "~31", ORANGE)
add_card(slide2, 10.1, 1.85, 2.5, 0.85, "TOTAL ABANDONMENTS", "29", RED)
add_card(slide2, 7.3, 2.85, 2.5, 0.85, "GRANT SUCCESS RATE", "98.7%", GREEN)
add_card(slide2, 10.1, 2.85, 2.5, 0.85, "STILL PENDING", "10", BLUE)

# Abandonment analysis
add_text(slide2, 7.3, 4.0, 5.5, 0.3, "ABANDONMENT ANALYSIS", size=11, color=RED, bold=True)
add_text(slide2, 7.3, 4.3, 5.5, 1.5,
    "Only 29 applications abandoned out of 3,000 (0.97%):\n"
    "  - 25 failed to respond to Office Actions\n"
    "  - 3 failed to pay the issue fee\n"
    "  - 1 abandoned after examiner's answer\n\n"
    "This is an exceptionally low abandonment rate. Most abandonments "
    "appear to be strategic decisions (letting non-essential designs lapse) "
    "rather than prosecution failures.",
    size=11, color=WHITE)

# Pipeline section at bottom
add_text(slide2, 0.5, 5.45, 12, 0.3, "ACTIVE EXAMINATION PIPELINE", size=11, color=BLUE, bold=True)

pipeline_data = [
    ("Ready for Examination", 2, BLUE),
    ("Office Action (Non-Final)", 3, ORANGE),
    ("Allowed / Notice of Allowance", 3, GREEN),
    ("Other Active", 2, PURPLE),
]

x = 0.5
for label, val, color in pipeline_data:
    add_card(slide2, x, 5.8, 2.8, 0.85, label.upper(), str(val), color)
    x += 3.05

add_text(slide2, 0.5, 6.85, 12, 0.3,
    "Only 10 applications currently in the pipeline — the vast majority of Apple's design portfolio is already granted.",
    size=11, color=DIMMED)

add_text(slide2, 0.5, 7.1, 12, 0.3, "Apple Inc.  |  Prosecution & Status Analysis  |  Data Source: USPTO PEDS", size=10, color=DIMMED, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3: Portfolio Analytics — Filing Trends, Inventors, Tech Areas
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide3)

add_text(slide3, 0.5, 0.3, 12, 0.6, "Portfolio Analytics", size=28, color=WHITE, bold=True)
add_text(slide3, 0.5, 0.8, 12, 0.3, "Filing Trends, Top Inventors & Law Firm Distribution", size=14, color=DIMMED)

shape = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.3), Pt(2))
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE
shape.line.fill.background()

# LEFT: Filing trends (bar chart via shapes)
add_text(slide3, 0.5, 1.45, 5.5, 0.3, "FILING TRENDS BY YEAR", size=11, color=BLUE, bold=True)

years_data = [
    ("2018", 327), ("2019", 417), ("2020", 484), ("2021", 448),
    ("2022", 502), ("2023", 403), ("2024", 321), ("2025", 98),
]
max_filings = 502
bar_base_y = 4.5
bar_max_h = 2.2
bar_w_each = 0.55
x_start = 0.7

for i, (year, count) in enumerate(years_data):
    bar_h = (count / max_filings) * bar_max_h
    x_pos = x_start + i * 0.72
    # Bar
    bar_shape = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x_pos), Inches(bar_base_y - bar_h),
        Inches(bar_w_each), Inches(bar_h)
    )
    bar_shape.fill.solid()
    bar_shape.fill.fore_color.rgb = BLUE
    bar_shape.line.fill.background()
    # Value on top
    add_text(slide3, x_pos - 0.05, bar_base_y - bar_h - 0.25, 0.65, 0.25, str(count), size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Year label
    add_text(slide3, x_pos - 0.05, bar_base_y + 0.05, 0.65, 0.25, year, size=9, color=DIMMED, align=PP_ALIGN.CENTER)

add_text(slide3, 0.5, 4.9, 6, 0.4, "Peak: 2022 (502)  |  Last 12 Months: 41 (-87% vs prior year)", size=10, color=DIMMED)

# RIGHT: Top inventors
add_text(slide3, 7.3, 1.45, 5.5, 0.3, "TOP 10 INVENTORS BY PATENT COUNT", size=11, color=BLUE, bold=True)

inventors = [
    ("Jonathan P. IVE", 849), ("Jody AKANA", 781), ("Richard P. HOWARTH", 736),
    ("Peter RUSSELL-CLARKE", 728), ("Bartley K. ANDRE", 722),
    ("Duncan Robert KERR", 721), ("M. Evans HANKEY", 704),
    ("Eugene Antony WHANG", 684), ("Jeremy BATAILLOU", 667),
    ("Shota AOYAGI", 657),
]
max_inv = 849
y = 1.85
for name, count in inventors:
    bar_w = (count / max_inv) * 3.5
    add_bar(slide3, 7.3, y, bar_w, 3.5, 0.2, YELLOW, name, str(count))
    y += 0.38

# Bottom section: Law firms + 12-month filing
add_text(slide3, 0.5, 5.4, 6, 0.3, "LAW FIRM DISTRIBUTION", size=11, color=BLUE, bold=True)

firms = [
    ("STERNE, KESSLER, GOLDSTEIN & FOX", 2452, BLUE),
    ("DesignLaw Group LLC", 428, GREEN),
    ("SAIDMAN DesignLaw Group", 119, YELLOW),
]
y = 5.75
for name, count, color in firms:
    bar_w = (count / 2452) * 5.5
    add_bar(slide3, 0.5, y, bar_w, 5.5, 0.22, color, name, f"{count:,}")
    y += 0.4

add_text(slide3, 7.3, 5.4, 5.5, 0.3, "FILING ACTIVITY", size=11, color=BLUE, bold=True)
add_card(slide3, 7.3, 5.75, 2.4, 0.85, "LAST 12 MONTHS", "41", RED)
add_card(slide3, 10.0, 5.75, 2.4, 0.85, "PRIOR 12 MONTHS", "308", GREEN)

add_text(slide3, 7.3, 6.7, 5, 0.5,
    "Significant year-over-year decline in design filings may indicate a strategic "
    "shift or consolidation phase in Apple's design IP approach.",
    size=10, color=DIMMED)

add_text(slide3, 0.5, 7.1, 12, 0.3, "Apple Inc.  |  Portfolio Analytics  |  1,680 Unique Inventors  |  3 Law Firms", size=10, color=DIMMED, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4: How It Was Built (Tech Stack — Layman Terms)
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide4)

add_text(slide4, 0.5, 0.3, 12, 0.6, "How It Was Built", size=28, color=WHITE, bold=True)
add_text(slide4, 0.5, 0.8, 12, 0.3, "Technology Stack & Architecture — In Plain English", size=14, color=DIMMED)

shape = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(12.3), Pt(2))
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE
shape.line.fill.background()

# Architecture flow (left to right)
add_text(slide4, 0.5, 1.45, 12, 0.3, "HOW DATA FLOWS THROUGH THE APP", size=11, color=BLUE, bold=True)

flow_items = [
    ("1. DATA SOURCE", "USPTO PEDS API\n(U.S. Patent Office)", "Raw patent data\nfor Apple Inc.", CYAN),
    ("2. DATA CLEANING", "clean_apple.py\n(Python Script)", "Parses names,\nformats dates,\nstandardizes titles", ORANGE),
    ("3. BACKEND API", "Flask Server\n(Python)", "Serves clean data\nto the browser\non port 5003", GREEN),
    ("4. FRONTEND APP", "React 19 + TypeScript\n(Web Browser)", "Interactive dashboard\nwith charts, tables,\nand search", BLUE),
]

x = 0.5
for title, tech, desc, color in flow_items:
    card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.85), Inches(2.8), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = color
    card.line.width = Pt(2)
    add_text(slide4, x + 0.15, 1.95, 2.5, 0.25, title, size=9, color=color, bold=True)
    add_text(slide4, x + 0.15, 2.2, 2.5, 0.5, tech, size=11, color=WHITE, bold=True)
    add_text(slide4, x + 0.15, 2.8, 2.5, 0.7, desc, size=10, color=DIMMED)
    # Arrow between cards
    if x < 9:
        arrow = slide4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.9), Inches(2.55), Inches(0.25), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = DIMMED
        arrow.line.fill.background()
    x += 3.15

# Tech stack table
add_text(slide4, 0.5, 4.0, 6, 0.3, "TECHNOLOGY STACK (WHAT EACH PIECE DOES)", size=11, color=BLUE, bold=True)

stack = [
    ("React 19", "The framework that builds the interactive web page — like the engine of a car", BLUE),
    ("TypeScript", "Adds type-safety to JavaScript — catches bugs before they happen", CYAN),
    ("Vite 7", "The development tool that makes the app load instantly during building", GREEN),
    ("Mantine v8", "Pre-built UI components (buttons, cards, layouts) — like building with LEGO", PURPLE),
    ("AG Grid v35", "The powerful data table — sort, filter, search 3,000 patents in real-time", YELLOW),
    ("Plotly.js", "The charting library — turns numbers into bar charts, pie charts, and funnels", ORANGE),
    ("Python Flask", "The backend server — reads the data file and sends it to the browser", GREEN),
    ("USPTO PEDS API", "The government data source — official patent records from the U.S. Patent Office", CYAN),
]

y = 4.35
for tech, desc, color in stack:
    # Color dot
    dot = slide4.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(y + 0.06), Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text(slide4, 0.85, y, 1.5, 0.25, tech, size=10, color=WHITE, bold=True)
    add_text(slide4, 2.4, y, 4.5, 0.25, desc, size=10, color=DIMMED)
    y += 0.3

# Right side: Why it matters
add_text(slide4, 7.3, 4.0, 5.5, 0.3, "WHY THIS MATTERS", size=11, color=GREEN, bold=True)
add_text(slide4, 7.3, 4.35, 5.3, 2.5,
    "This tool was built entirely with free, open-source technologies — "
    "no expensive licenses or subscriptions required.\n\n"
    "Commercial patent analytics platforms (like PatSnap, Anaqua, "
    "or Darts-ip) cost $50,000 to $200,000+ per year.\n\n"
    "This dashboard delivers the same core functionality:\n"
    "  - Portfolio overview with key health metrics\n"
    "  - Prosecution pipeline tracking (OAs, grants, abandonments)\n"
    "  - Filing trend analysis over time\n"
    "  - Inventor intelligence and collaboration patterns\n"
    "  - Searchable patent database with CSV export\n"
    "  - Maintenance fee deadline tracking\n\n"
    "All running in a browser, built in hours, free to maintain.",
    size=11, color=WHITE)

add_text(slide4, 0.5, 7.1, 12, 0.3, "Apple Inc.  |  Built with React 19 + Flask + Plotly  |  Open-Source Stack  |  Zero License Cost", size=10, color=DIMMED, align=PP_ALIGN.CENTER)


# Save
prs.save("/Users/fgomez/Developer/Apple/Apple_Patent_Explorer_Executive_Summary.pptx")
print("PowerPoint saved: Apple_Patent_Explorer_Executive_Summary.pptx")
